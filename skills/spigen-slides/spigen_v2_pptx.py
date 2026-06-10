"""
spigen_v2_pptx.py — Spigen Slides V2 PPTX builder

V2 output path:
1) Build a local editable PPTX at 1920 × 1080pt.
2) Convert it to native Google Slides with gws Drive upload:
   gws drive files create --json '{"name":"...", "mimeType":"application/vnd.google-apps.presentation"}' \
     --upload deck.pptx \
     --upload-content-type application/vnd.openxmlformats-officedocument.presentationml.presentation

Why PPTX-first:
- Native Google Slides templates often keep a 720 × 405pt pageSize.
- True 1920 × 1080pt output requires creating/importing a deck with that pageSize.
"""

from __future__ import annotations

import json
import subprocess
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt


W_PT = 1920
H_PT = 1080
MARGIN_PT = 96
CONTENT_TOP_PT = 300
CONTENT_BOTTOM_PT = 1016
CORNER_RADIUS_PT = 4  # target radius; PPTX primitive uses crisp rectangles by default.

BG = RGBColor(0, 0, 0)
SURFACE = RGBColor(14, 14, 14)
SURFACE_2 = RGBColor(24, 24, 24)
TEXT = RGBColor(245, 245, 245)
DIM = RGBColor(172, 172, 172)
FAINT = RGBColor(105, 105, 105)
LINE = RGBColor(50, 50, 50)
ORANGE = RGBColor(255, 107, 26)
GREEN = RGBColor(52, 168, 83)
YELLOW = RGBColor(245, 176, 65)

FONT_KR = "Noto Sans KR"
FONT_EN = "Proxima Nova"


def _pt(v):
    return Pt(v)


class SpigenV2Deck:
    """Minimal 1920×1080pt dark-guide builder for editable PPTX → Google Slides."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Pt(W_PT)
        self.prs.slide_height = Pt(H_PT)
        self.blank = self.prs.slide_layouts[6]

    def slide(self):
        slide = self.prs.slides.add_slide(self.blank)
        self.rect(slide, 0, 0, W_PT, H_PT, fill=BG, outline=BG, crisp=True)
        return slide

    def text(self, slide, text, x, y, w, h, size=28, color=TEXT, bold=False,
             font=FONT_KR, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line=1.0):
        box = slide.shapes.add_textbox(_pt(x), _pt(y), _pt(w), _pt(h))
        tf = box.text_frame
        tf.clear()
        tf.margin_left = _pt(0)
        tf.margin_right = _pt(0)
        tf.margin_top = _pt(0)
        tf.margin_bottom = _pt(0)
        tf.vertical_anchor = valign
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.line_spacing = line
        for run in p.runs:
            run.font.name = font
            run.font.size = _pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
        return box

    def rect(self, slide, x, y, w, h, fill=SURFACE, outline=LINE, crisp=True):
        # Use crisp rectangles by default. Avoid large rounded cards / AI-template look.
        shape_type = MSO_SHAPE.RECTANGLE if crisp else MSO_SHAPE.ROUNDED_RECTANGLE
        shp = slide.shapes.add_shape(shape_type, _pt(x), _pt(y), _pt(w), _pt(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = outline
        shp.line.width = _pt(1)
        return shp

    def line(self, slide, x, y, w, h=2, color=LINE):
        return self.rect(slide, x, y, w, h, fill=color, outline=color, crisp=True)

    def vline(self, slide, x, y, h, color=LINE):
        return self.line(slide, x, y, 2, h, color)

    def dot(self, slide, x, y, color=ORANGE, d=13):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, _pt(x), _pt(y), _pt(d), _pt(d))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        return shp

    def arrow(self, slide, x1, y1, x2, y2, color=ORANGE):
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _pt(x1), _pt(y1), _pt(x2), _pt(y2))
        conn.line.color.rgb = color
        conn.line.width = _pt(3.5)
        conn.line.end_arrowhead = True
        return conn

    def pill(self, slide, text, x, y, w, color=ORANGE, text_color=BG):
        self.rect(slide, x, y, w, 40, fill=color, outline=color, crisp=True)
        self.text(slide, text, x, y + 10, w, 18, 16, text_color, True, align=PP_ALIGN.CENTER)

    def header(self, slide, eyebrow, title, lead=None, page=None):
        self.text(slide, eyebrow.upper(), 96, 72, 500, 28, 17, ORANGE, True, FONT_EN)
        self.text(slide, title, 96, 116, 1320, 68, 50, TEXT, True)
        if lead:
            self.text(slide, lead, 98, 198, 1300, 40, 23, DIM)
        if page:
            self.text(slide, f"{page:02d}", 1770, 72, 70, 24, 16, FAINT, True, FONT_EN, PP_ALIGN.RIGHT)

    def cover(self, title, subtitle="", meta="디자인부문ㅣ패키지디자인팀\n한원진 담당", date=None):
        """Create the preserved Spigen orange cover.

        V2 changes content-slide design only. The cover keeps the existing
        orange template look: orange full background, black title/meta/date,
        and a simple editable black logo mark at top-right.
        """
        s = self.prs.slides.add_slide(self.blank)
        self.rect(s, 0, 0, W_PT, H_PT, fill=ORANGE, outline=ORANGE, crisp=True)
        self.text(s, title, 96, 120, 900, 220, 76, BG, True)
        if subtitle:
            self.text(s, subtitle, 100, 380, 900, 38, 28, BG)
        self.text(s, meta, 100, 865, 620, 90, 23, BG)
        self.text(s, date or datetime.now().strftime("%Y.%m.%d"), 1625, 890, 230, 34, 23, BG, align=PP_ALIGN.RIGHT)
        for x, y in [(1728, 150), (1778, 116), (1778, 184), (1828, 150)]:
            self.dot(s, x, y, BG, 34)
        for x1, y1, x2, y2 in [(1756, 142, 1792, 124), (1756, 158, 1792, 176), (1805, 150, 1818, 150)]:
            conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(x1), Pt(y1), Pt(x2), Pt(y2))
            conn.line.color.rgb = BG
            conn.line.width = Pt(8)
        return s

    def save(self, path):
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return str(path)

    @staticmethod
    def import_to_google_slides(path, title):
        body = {"name": title, "mimeType": "application/vnd.google-apps.presentation"}
        result = subprocess.run(
            [
                "gws", "drive", "files", "create",
                "--json", json.dumps(body, ensure_ascii=False),
                "--upload", str(path),
                "--upload-content-type",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            raise RuntimeError((result.stdout + "\n" + result.stderr).strip())
        return json.loads(result.stdout)
